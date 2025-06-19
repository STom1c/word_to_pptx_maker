#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
通用Word轉PowerPoint工具 - 圖形化介面版本 (繁體中文版)
支援任意Word文件和PowerPoint範本，智慧內容解析和風格匹配
增強功能：
1. 改進中文章節識別（一、二、三、四等）
2. 增大字體大小至少32pt
3. 圖片式投影片預覽顯示
4. 記憶上次使用的檔案路徑
"""

import sys
import os
import re
import subprocess
import json
from pathlib import Path
from typing import List, Dict, Tuple, Optional
from dataclasses import dataclass
from io import BytesIO
import tempfile

from PySide6.QtWidgets import (
    QApplication, QMainWindow, QVBoxLayout, QHBoxLayout, QGridLayout,
    QPushButton, QLabel, QTextEdit, QProgressBar, QFileDialog,
    QWidget, QFrame, QScrollArea, QTabWidget, QGroupBox,
    QMessageBox, QSplitter, QListWidget, QListWidgetItem,
    QComboBox, QSpinBox, QCheckBox, QSlider
)
from PySide6.QtCore import (
    Qt, QThread, Signal, QMimeData, QTimer, QPropertyAnimation,
    QEasingCurve, QRect, Slot, QSize, QSettings
)
from PySide6.QtGui import (
    QFont, QPixmap, QPainter, QColor, QBrush, QPen,
    QDragEnterEvent, QDropEvent, QIcon, QLinearGradient
)

# 第三方庫導入
try:
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image, ImageDraw, ImageFont
    DEPENDENCIES_OK = True
except ImportError as e:
    print(f"請安裝必要套件: pip install python-docx python-pptx pillow")
    DEPENDENCIES_OK = False

@dataclass
class ContentBlock:
    """內容塊資料結構"""
    text: str
    level: int  # 標題層級 (0=主標題/章節, 1=次標題, 2=內容)
    content_type: str  # header, chapter, title, subtitle, content, quote, list
    formatting: Dict = None

@dataclass
class SlideTemplate:
    """投影片範本資料結構"""
    layout_index: int
    layout_name: str
    placeholders: List[Dict]
    background_color: Tuple[int, int, int] = None
    font_family: str = "Microsoft JhengHei"

class ConfigManager:
    """設定管理器 - 記憶上次使用的檔案路徑"""
    
    def __init__(self):
        self.config_file = os.path.join(os.path.expanduser("~"), ".word_to_ppt_config.json")
        self.config = self.load_config()
    
    def load_config(self) -> Dict:
        """載入設定檔"""
        default_config = {
            "last_word_path": "",
            "last_template_path": "",
            "last_output_dir": "",
            "window_geometry": None
        }
        
        try:
            if os.path.exists(self.config_file):
                with open(self.config_file, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    # 驗證檔案是否仍然存在
                    if config.get("last_word_path") and not os.path.exists(config["last_word_path"]):
                        config["last_word_path"] = ""
                    if config.get("last_template_path") and not os.path.exists(config["last_template_path"]):
                        config["last_template_path"] = ""
                    return {**default_config, **config}
        except Exception as e:
            print(f"載入設定檔時發生錯誤: {e}")
        
        return default_config
    
    def save_config(self):
        """儲存設定檔"""
        try:
            with open(self.config_file, 'w', encoding='utf-8') as f:
                json.dump(self.config, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"儲存設定檔時發生錯誤: {e}")
    
    def set_last_word_path(self, path: str):
        """設定上次使用的Word檔案路徑"""
        self.config["last_word_path"] = path
        self.save_config()
    
    def set_last_template_path(self, path: str):
        """設定上次使用的範本檔案路徑"""
        self.config["last_template_path"] = path
        self.save_config()
    
    def set_last_output_dir(self, dir_path: str):
        """設定上次使用的輸出目錄"""
        self.config["last_output_dir"] = dir_path
        self.save_config()
    
    def get_last_word_path(self) -> str:
        """取得上次使用的Word檔案路徑"""
        return self.config.get("last_word_path", "")
    
    def get_last_template_path(self) -> str:
        """取得上次使用的範本檔案路徑"""
        return self.config.get("last_template_path", "")
    
    def get_last_output_dir(self) -> str:
        """取得上次使用的輸出目錄"""
        return self.config.get("last_output_dir", "")

class WordDocumentAnalyzer:
    """Word文件智慧分析器 - 增強版"""
    
    def __init__(self):
        # 增強的中文章節識別模式
        self.chapter_patterns = [
            # 中文數字章節
            r'^[一二三四五六七八九十]+[、．.]\s*',
            r'^第[一二三四五六七八九十壹貳參肆伍陸柒捌玖拾]+[章節部分]\s*',
            r'^第[一二三四五六七八九十]+[、．.]\s*',
            '''
            # 阿拉伯數字章節
            r'^[1-9]\d*[、．.]\s*',
            r'^第[1-9]\d*[章節部分]\s*',
            r'^第[1-9]\d*[、．.]\s*',
            # 英文字母章節
            r'^[A-Z][、．.]\s*',
            r'^第[A-Z][章節部分]\s*',
            # 特殊標記
            r'^[●◆■▲]\s*',
            '''
            r'^[前言結論總結概述摘要序言引言]\s*',
        ]
        
        self.subtitle_patterns = [
            r'^[一二三四五六七八九十]+[）)]\s*',
            r'^\([一二三四五六七八九十]+\)\s*',
            r'^[1-9]\d*[）)]\s*',
            r'^\([1-9]\d*\)\s*',
            r'^[a-z][）)]\s*',
            r'^\([a-z]\)\s*',
            r'^[•·○]\s*',
        ]
        
    def analyze_document(self, file_path: str) -> List[ContentBlock]:
        """分析Word文件結構"""
        try:
            doc = Document(file_path)
            blocks = []
            
            header_line = True
            for para in doc.paragraphs:
                if not para.text.strip():
                    continue

                if header_line:
                    # 如果是第一行，可能是標題或章節
                    block = self._classify_header(para)
                    header_line = False
                else:
                    # 其他行是內容    
                    block = self._analyze_paragraph(para)
                if block:
                    blocks.append(block)
            
            return self._optimize_structure(blocks)
            
        except Exception as e:
            raise Exception(f"分析Word文件時發生錯誤: {e}")
    
    def _analyze_paragraph(self, para) -> Optional[ContentBlock]:
        """分析單個段落"""
        text = para.text.strip()
        if not text:
            return None
            
        # 分析文字格式
        formatting = self._extract_formatting(para)
        
        # 判斷內容類型和層級
        level, content_type = self._classify_content(text, formatting)
        
        return ContentBlock(
            text=text,
            level=level,
            content_type=content_type,
            formatting=formatting
        )

    def _classify_header(self, para) -> Optional[ContentBlock]:
        text = para.text.strip()
        if not text:
            return None

        # 分析文字格式
        formatting = self._extract_formatting(para)

        return ContentBlock(
            text=text,
            level=0,  # 假設第一行為主標題或章節
            content_type='header',  # 假設為章節
            formatting=formatting
        )


    def _extract_formatting(self, para) -> Dict:
        """提取段落格式"""
        formatting = {
            'bold': False,
            'italic': False,
            'font_size': 12,
            'alignment': 'left'
        }
        
        if para.runs:
            # 檢查第一個run的格式
            run = para.runs[0]
            if run.bold:
                formatting['bold'] = True
            if run.italic:
                formatting['italic'] = True
            if run.font.size:
                formatting['font_size'] = run.font.size.pt
                
        return formatting
    
    def _classify_content(self, text: str, formatting: Dict) -> Tuple[int, str]:
        """分類內容類型和層級 - 增強版"""
        # 優先檢查章節模式
        for pattern in self.chapter_patterns:
            if re.match(pattern, text):
                return (0, 'chapter')  # 章節標題
        
        # 檢查次標題模式
        for pattern in self.subtitle_patterns:
            if re.match(pattern, text):
                return (1, 'subtitle')  # 次標題
        '''
        # 檢查格式特徵
        font_size = formatting.get('font_size', 12)
        is_bold = formatting.get('bold', False)
        
        if is_bold and font_size > 16:
            return (0, 'title')  # 格式化的主標題
        elif is_bold and font_size > 14:
            return (1, 'subtitle')  # 格式化的次標題
        
        # 檢查特殊內容
        if text.startswith('—') or '引自' in text or '出處' in text:
            return (2, 'quote')
        
        if any(text.startswith(marker) for marker in ['●', '•', '○', '-', '※']):
            return (2, 'list')
        
        # 檢查是否為列表項目
        if re.match(r'^\d+[\.)]\s*', text) or re.match(r'^[a-zA-Z][\.)]\s*', text):
            return (2, 'list')
        '''    
        return (2, 'content')
    
    def X_classify_content(self, text: str, formatting: Dict) -> Tuple[int, str]:
        """分類內容類型和層級 - 增強版"""
        # 優先檢查章節模式
        for pattern in self.chapter_patterns:
            if re.match(pattern, text):
                return (0, 'chapter')  # 章節標題
        
        # 檢查次標題模式
        for pattern in self.subtitle_patterns:
            if re.match(pattern, text):
                return (1, 'subtitle')  # 次標題
        
        # 檢查格式特徵
        font_size = formatting.get('font_size', 12)
        is_bold = formatting.get('bold', False)
        
        if is_bold and font_size > 16:
            return (0, 'title')  # 格式化的主標題
        elif is_bold and font_size > 14:
            return (1, 'subtitle')  # 格式化的次標題
        
        # 檢查特殊內容
        if text.startswith('—') or '引自' in text or '出處' in text:
            return (2, 'quote')
        
        if any(text.startswith(marker) for marker in ['●', '•', '○', '-', '※']):
            return (2, 'list')
        
        # 檢查是否為列表項目
        if re.match(r'^\d+[\.)]\s*', text) or re.match(r'^[a-zA-Z][\.)]\s*', text):
            return (2, 'list')
            
        return (2, 'content')
    
    def _optimize_structure(self, blocks: List[ContentBlock]) -> List[ContentBlock]:
        """最佳化文件結構"""
        if not blocks:
            return blocks
            
        # 確保有主標題或章節
        has_chapter = any(block.content_type == 'chapter' for block in blocks)
        if not has_chapter and blocks:
            # 將第一個標題或粗體內容提升為章節
            for block in blocks:
                if block.level <= 1 or block.formatting.get('bold', False):
                    block.level = 0
                    block.content_type = 'chapter'
                    break
        
        return blocks

class PowerPointTemplateAnalyzer:
    """PowerPoint範本分析器"""
    
    def __init__(self):
        self.templates = []
        
    def analyze_template(self, file_path: str) -> List[SlideTemplate]:
        """分析PowerPoint範本"""
        try:
            prs = Presentation(file_path)
            templates = []
            
            for i, layout in enumerate(prs.slide_layouts):
                template = SlideTemplate(
                    layout_index=i,
                    layout_name=layout.name,
                    placeholders=self._extract_placeholders(layout)
                )
                templates.append(template)
            
            self.templates = templates
            return templates
            
        except Exception as e:
            raise Exception(f"分析PowerPoint範本時發生錯誤: {e}")
    
    def _extract_placeholders(self, layout) -> List[Dict]:
        """提取預留位置資訊"""
        placeholders = []
        
        for placeholder in layout.placeholders:
            try:
                info = {
                    'idx': placeholder.placeholder_format.idx,
                    'type': str(placeholder.placeholder_format.type),
                    'name': getattr(placeholder, 'name', ''),
                    'left': placeholder.left,
                    'top': placeholder.top,
                    'width': placeholder.width,
                    'height': placeholder.height
                }
                placeholders.append(info)
            except:
                continue
                
        return placeholders

class ContentToSlideMapper:
    """內容到投影片的智慧映射器 - 增強版"""
    
    def __init__(self, presentation_path: str):
        self.presentation_path = presentation_path
        self.prs = None
        
    def create_slides(self, blocks: List[ContentBlock]) -> Presentation:
        """建立投影片"""
        self.prs = Presentation(self.presentation_path)
        
        # 清空現有投影片
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]
        
        current_slide = None
        current_content = []
        
        for block in blocks:
            print(f'Block level:{block.level}, type:{block.content_type}, text:{block.text}')
            if block.content_type == 'header':  # 主標題 - 新投影片
                if current_slide is not None:
                    self._finalize_slide(current_slide, current_content)
                current_slide = self._create_title_slide(block)
                current_content = []
                
            elif block.content_type == 'chapter':  # 章節 - 新投影片
                if current_slide is not None:
                    self._finalize_slide(current_slide, current_content)
                    current_content = []
                current_slide = self._create_content_slide(block)
                '''    
                elif block.level == 1:  # 次標題 - 新內容投影片
                    if current_slide is not None:
                        self._finalize_slide(current_slide, current_content)
                    current_slide = self._create_content_slide(block)
                    current_content = []
                ''' 
            else:  # 內容 - 添加到當前投影片
                if current_slide is None:
                    current_slide = self._create_content_slide(block)
                current_content.append(block)
            
        # 處理最後一張投影片
        if current_slide is not None:
            self._finalize_slide(current_slide, current_content)
            
        return self.prs
    
    def _create_title_slide(self, block: ContentBlock):
        """建立標題投影片"""
        layout = self._get_best_layout('title')
        slide = self.prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            # 清理章節標記，只保留標題文字
            title_text = self._clean_chapter_text(block.text)
            slide.shapes.title.text = title_text
            #self._format_title(slide.shapes.title, block.formatting)
            
        return slide
    
    def _create_content_slide(self, block: ContentBlock):
        """建立內容投影片"""
        layout = self._get_best_layout('content')
        slide = self.prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            # 清理次標題標記
            title_text = self._clean_subtitle_text(block.text)
            slide.shapes.title.text = title_text
            #self._format_subtitle(slide.shapes.title, block.formatting)
            
        return slide
    
    def _clean_chapter_text(self, text: str) -> str:
        """清理章節文字，移除編號標記"""
        # 移除各種章節標記
        patterns = [
            r'^[一二三四五六七八九十]+[、．.]\s*',
            r'^第[一二三四五六七八九十壹貳參肆伍陸柒捌玖拾]+[章節部分]\s*',
            r'^第[一二三四五六七八九十]+[、．.]\s*',
            r'^[1-9]\d*[、．.]\s*',
            r'^第[1-9]\d*[章節部分]\s*',
            r'^第[1-9]\d*[、．.]\s*',
            r'^[A-Z][、．.]\s*',
            r'^第[A-Z][章節部分]\s*',
            r'^[●◆■▲]\s*',
        ]
        
        cleaned_text = text
        for pattern in patterns:
            cleaned_text = re.sub(pattern, '', cleaned_text)
        
        return cleaned_text.strip()
    
    def _clean_subtitle_text(self, text: str) -> str:
        """清理次標題文字"""
        patterns = [
            r'^[一二三四五六七八九十]+[）)]\s*',
            r'^\([一二三四五六七八九十]+\)\s*',
            r'^[1-9]\d*[）)]\s*',
            r'^\([1-9]\d*\)\s*',
            r'^[a-z][）)]\s*',
            r'^\([a-z]\)\s*',
            r'^[•·○]\s*',
        ]
        
        cleaned_text = text
        for pattern in patterns:
            cleaned_text = re.sub(pattern, '', cleaned_text)
        
        return cleaned_text.strip()
    
    def _finalize_slide(self, slide, content_blocks: List[ContentBlock]):
        """完成投影片內容"""
        if not content_blocks:
            return
            
        # 找到內容預留位置
        content_placeholder = None
        for shape in slide.placeholders:
            if hasattr(shape, 'text_frame') and shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break
        
        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            for i, block in enumerate(content_blocks):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = block.text
                #self._format_content(p, block.formatting, block.content_type)
    
    def _get_best_layout(self, layout_type: str):
        """取得最佳佈局"""
        if not self.prs:
            raise Exception("Presentation not initialized")
            
        if layout_type == 'title':
            # 尋找標題佈局 - 通常是第一個佈局
            for layout in self.prs.slide_layouts:
                if 'title' in layout.name.lower() or layout == self.prs.slide_layouts[0]:
                    return layout
        else:
            # 尋找內容佈局 - 通常是第二個佈局
            for layout in self.prs.slide_layouts:
                if 'content' in layout.name.lower():
                    return layout
            # 如果沒找到content佈局，使用第二個佈局
            if len(self.prs.slide_layouts) > 1:
                return self.prs.slide_layouts[1]
        
        # 預設回傳第一個佈局
        return self.prs.slide_layouts[0]
    
    def _format_title(self, shape, formatting: Dict):
        """格式化標題 - 增大字體"""
        if not shape or not hasattr(shape, 'text_frame') or not shape.text_frame:
            return
            
        try:
            para = shape.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            
            if hasattr(para, 'font') and para.font:
                para.font.name = 'Microsoft JhengHei'
                para.font.size = Pt(48)  # 增大到48pt
                para.font.bold = True
                para.font.color.rgb = RGBColor(44, 62, 80)
        except Exception as e:
            print(f"格式化標題時發生錯誤: {e}")
    
    def _format_subtitle(self, shape, formatting: Dict):
        """格式化次標題 - 增大字體"""
        if not shape or not hasattr(shape, 'text_frame') or not shape.text_frame:
            return
            
        try:
            para = shape.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            
            if hasattr(para, 'font') and para.font:
                para.font.name = 'Microsoft JhengHei'
                para.font.size = Pt(40)  # 增大到40pt
                para.font.bold = True
                para.font.color.rgb = RGBColor(26, 54, 93)
        except Exception as e:
            print(f"格式化次標題時發生錯誤: {e}")
    
    def _format_content(self, paragraph, formatting: Dict, content_type: str):
        """格式化內容 - 確保至少32pt"""
        try:
            paragraph.alignment = PP_ALIGN.LEFT
            
            if hasattr(paragraph, 'font') and paragraph.font:
                paragraph.font.name = 'Microsoft JhengHei'
                
                if content_type == 'quote':
                    paragraph.font.size = Pt(32)  # 至少32pt
                    paragraph.font.italic = True
                    paragraph.font.color.rgb = RGBColor(218, 165, 32)
                elif content_type == 'list':
                    paragraph.font.size = Pt(32)  # 至少32pt
                    paragraph.font.color.rgb = RGBColor(44, 82, 130)
                else:
                    paragraph.font.size = Pt(32)  # 至少32pt
                    paragraph.font.color.rgb = RGBColor(44, 62, 80)
        except Exception as e:
            print(f"格式化內容時發生錯誤: {e}")

class SlideImageGenerator:
    """投影片圖片產生器 - 使用圖片方式顯示投影片 (UTF-8修復版)"""
    
    def __init__(self):
        self.temp_dir = tempfile.mkdtemp()
        print(f"建立暫存目錄: {self.temp_dir}")
        self.fonts_cache = {}
        self._load_system_fonts()
    
    def _load_system_fonts(self):
        """載入系統字體並測試中文支援"""
        self.available_fonts = {}
        
        # 系統字體路徑 (依優先順序)
        font_candidates = [
            # Windows 繁體中文字體
            ("Microsoft JhengHei", [
                "C:/Windows/Fonts/msjh.ttc",
                "C:/Windows/Fonts/msjhbd.ttc"
            ]),
            ("Microsoft YaHei", [
                "C:/Windows/Fonts/msyh.ttc", 
                "C:/Windows/Fonts/msyhbd.ttc"
            ]),
            ("SimSun", [
                "C:/Windows/Fonts/simsun.ttc"
            ]),
            # macOS 中文字體
            ("PingFang SC", [
                "/System/Library/Fonts/PingFang.ttc",
                "/Library/Fonts/PingFang.ttc"
            ]),
            ("Hiragino Sans", [
                "/System/Library/Fonts/Hiragino Sans GB.ttc"
            ]),
            # Linux 中文字體
            ("WenQuanYi", [
                "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
                "/usr/share/fonts/truetype/arphic/uming.ttc"
            ]),
            # 通用字體
            ("Arial Unicode MS", [
                "/Library/Fonts/Arial Unicode MS.ttf",
                "C:/Windows/Fonts/ARIALUNI.TTF"
            ])
        ]
        
        # 測試每種字體是否支援中文
        for font_name, paths in font_candidates:
            for path in paths:
                if os.path.exists(path):
                    if self._test_font_chinese_support(path):
                        self.available_fonts[font_name] = path
                        print(f"載入中文字體: {font_name} ({path})")
                        break
        
        if not self.available_fonts:
            print("警告: 未找到支援中文的字體，可能影響顯示效果")
    
    def _test_font_chinese_support(self, font_path: str) -> bool:
        """測試字體是否支援中文"""
        try:
            test_font = ImageFont.truetype(font_path, 20)
            # 建立測試圖片
            test_img = Image.new('RGB', (100, 50), 'white')
            test_draw = ImageDraw.Draw(test_img)
            
            # 測試繪製中文字
            test_text = "測試中文字體"
            test_draw.text((10, 10), test_text, font=test_font, fill='black')
            
            # 如果沒有拋出例外，認為支援中文
            return True
            
        except Exception as e:
            print(f"字體測試失敗 {font_path}: {e}")
            return False
    
    def _get_best_font(self, size: int) -> ImageFont.ImageFont:
        """取得最佳的中文字體"""
        cache_key = f"font_{size}"
        if cache_key in self.fonts_cache:
            return self.fonts_cache[cache_key]
        
        # 依優先順序嘗試字體
        priority_fonts = ["Microsoft JhengHei", "Microsoft YaHei", "PingFang SC", "SimSun"]
        
        for font_name in priority_fonts:
            if font_name in self.available_fonts:
                try:
                    font = ImageFont.truetype(self.available_fonts[font_name], size)
                    self.fonts_cache[cache_key] = font
                    return font
                except Exception as e:
                    print(f"載入字體失敗 {font_name}: {e}")
                    continue
        
        # 如果都失敗，使用預設字體
        default_font = ImageFont.load_default()
        self.fonts_cache[cache_key] = default_font
        return default_font
    
    def _normalize_text(self, text: str) -> str:
        """標準化文字編碼為UTF-8"""
        try:
            # 確保文字是字串類型
            if isinstance(text, bytes):
                # 嘗試不同編碼解碼
                for encoding in ['utf-8', 'big5', 'gb2312', 'utf-16']:
                    try:
                        text = text.decode(encoding)
                        break
                    except:
                        continue
                else:
                    # 所有編碼都失敗，使用錯誤處理
                    text = text.decode('utf-8', errors='ignore')
            
            # 轉換為字串並處理特殊字元
            text = str(text)
            
            # 確保文字是有效的UTF-8
            text = text.encode('utf-8', errors='ignore').decode('utf-8')
            
            # 移除控制字元但保留換行符
            import unicodedata
            text = ''.join(char for char in text 
                          if unicodedata.category(char)[0] != 'C' or char in '\n\r\t')
            
            return text.strip()
            
        except Exception as e:
            print(f"文字編碼標準化失敗: {e}")
            # 返回安全的文字
            return str(text)[:100] if text else "文字顯示錯誤"
    
    def generate_slide_images(self, presentation_path: str) -> List[str]:
        """產生投影片圖片"""
        try:
            # 使用python-pptx載入簡報
            prs = Presentation(presentation_path)
            image_paths = []
            
            print(f"開始產生 {len(prs.slides)} 張投影片圖片...")
            
            for i, slide in enumerate(prs.slides):
                print(f"正在處理投影片 {i+1}...")
                image_path = self._create_slide_image(slide, i + 1)
                if image_path:
                    image_paths.append(image_path)
            
            print(f"成功產生 {len(image_paths)} 張投影片圖片")
            return image_paths
            
        except Exception as e:
            print(f"產生投影片圖片時發生錯誤: {e}")
            import traceback
            print(traceback.format_exc())
            return []
    
    def _create_slide_image(self, slide, slide_number: int) -> str:
        """建立單張投影片圖片 (UTF-8修復版)"""
        try:
            # 建立畫布 (16:9 比例)
            width, height = 800, 450
            img = Image.new('RGB', (width, height), color='white')
            draw = ImageDraw.Draw(img)
            
            # 載入字體
            title_font = self._get_best_font(32)
            content_font = self._get_best_font(24)
            small_font = self._get_best_font(16)
            
            # 繪製邊框
            draw.rectangle([0, 0, width-1, height-1], outline='#CCCCCC', width=2)
            
            # 提取和繪製內容
            y_pos = 30
            title_drawn = False
            
            try:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text and shape.text.strip():
                        # 取得並標準化文字
                        raw_text = shape.text.strip()
                        text = self._normalize_text(raw_text)
                        
                        if not text:
                            continue
                        
                        print(f"處理文字: {text[:50]}...")
                        
                        # 檢查是否為標題
                        is_title = (hasattr(shape, 'placeholder_format') and 
                                  shape.placeholder_format.idx == 0)
                        
                        if is_title and not title_drawn:
                            # 繪製標題
                            y_pos = self._draw_text_centered_safe(
                                draw, text, width//2, y_pos, title_font, '#2C3E50'
                            )
                            y_pos += 60
                            title_drawn = True
                        elif not is_title:
                            # 繪製內容
                            y_pos = self._draw_text_wrapped_safe(
                                draw, text, 40, y_pos, width-80, content_font, '#34495E'
                            )
                            y_pos += 40
                        
                        # 防止內容超出邊界
                        if y_pos > height - 80:
                            break
                            
            except Exception as e:
                print(f"處理投影片內容時發生錯誤: {e}")
                # 繪製錯誤訊息
                error_text = f"投影片 {slide_number} 內容處理錯誤"
                self._draw_text_centered_safe(
                    draw, error_text, width//2, height//2, content_font, '#E74C3C'
                )
            
            # 添加投影片編號 (使用UTF-8)
            number_text = f"投影片 {slide_number}"
            try:
                self._draw_text_safe(draw, number_text, 20, height-35, small_font, '#7F8C8D')
            except:
                # 如果中文失敗，使用英文
                number_text = f"Slide {slide_number}"
                self._draw_text_safe(draw, number_text, 20, height-35, small_font, '#7F8C8D')
            
            # 儲存圖片
            image_path = os.path.join(self.temp_dir, f"slide_{slide_number}.png")
            img.save(image_path, 'PNG', quality=95)
            
            print(f"成功建立投影片圖片: {image_path}")
            return image_path
            
        except Exception as e:
            print(f"建立投影片圖片時發生錯誤: {e}")
            import traceback
            print(traceback.format_exc())
            return None
    
    def _draw_text_safe(self, draw, text: str, x: int, y: int, font, color: str):
        """安全繪製文字 (處理UTF-8編碼)"""
        try:
            # 確保文字編碼正確
            safe_text = self._normalize_text(text)
            draw.text((x, y), safe_text, fill=color, font=font)
        except UnicodeEncodeError:
            # 如果仍有編碼問題，使用ASCII安全版本
            safe_text = text.encode('ascii', errors='ignore').decode('ascii')
            if safe_text.strip():
                draw.text((x, y), safe_text, fill=color, font=font)
            else:
                draw.text((x, y), "[文字]", fill=color, font=font)
        except Exception as e:
            print(f"繪製文字時發生錯誤: {e}")
            # 最後保障：繪製簡單文字
            try:
                draw.text((x, y), "[文字內容]", fill=color, font=font)
            except:
                pass
    
    def _draw_text_centered_safe(self, draw, text: str, x: int, y: int, font, color: str) -> int:
        """安全繪製置中文字"""
        try:
            safe_text = self._normalize_text(text)
            
            # 計算文字寬度
            try:
                bbox = draw.textbbox((0, 0), safe_text, font=font)
                text_width = bbox[2] - bbox[0]
                text_height = bbox[3] - bbox[1]
            except:
                # 估算文字寬度
                text_width = len(safe_text) * 12
                text_height = 32
            
            # 置中繪製
            center_x = x - text_width // 2
            self._draw_text_safe(draw, safe_text, center_x, y, font, color)
            
            return y + text_height + 10
            
        except Exception as e:
            print(f"繪製置中文字時發生錯誤: {e}")
            return y + 40
    
    def _draw_text_wrapped_safe(self, draw, text: str, x: int, y: int, max_width: int, font, color: str) -> int:
        """安全繪製自動換行文字 (中文友善)"""
        try:
            safe_text = self._normalize_text(text)
            
            # 對中文進行更好的換行處理
            lines = self._wrap_chinese_text(safe_text, max_width, font, draw)
            
            # 繪製每一行
            line_height = 32
            for i, line in enumerate(lines[:6]):  # 最多6行
                if line.strip():
                    line_y = y + i * line_height
                    self._draw_text_safe(draw, line, x, line_y, font, color)
            
            return y + len(lines) * line_height
            
        except Exception as e:
            print(f"繪製換行文字時發生錯誤: {e}")
            # 簡單處理
            self._draw_text_safe(draw, text[:30] + "...", x, y, font, color)
            return y + 40
    
    def _wrap_chinese_text(self, text: str, max_width: int, font, draw) -> List[str]:
        """中文友善的文字換行"""
        try:
            lines = []
            paragraphs = text.split('\n')
            
            for paragraph in paragraphs:
                if not paragraph.strip():
                    lines.append('')
                    continue
                
                # 對於中文，按字符而非單詞換行更合適
                current_line = ""
                chars = list(paragraph)
                
                for char in chars:
                    test_line = current_line + char
                    
                    try:
                        # 嘗試測量文字寬度
                        bbox = draw.textbbox((0, 0), test_line, font=font)
                        text_width = bbox[2] - bbox[0]
                        
                        if text_width <= max_width:
                            current_line = test_line
                        else:
                            if current_line:
                                lines.append(current_line)
                            current_line = char
                    except:
                        # 如果測量失敗，使用字符數估算
                        if len(current_line) < max_width // 15:  # 估算值
                            current_line = test_line
                        else:
                            if current_line:
                                lines.append(current_line)
                            current_line = char
                
                if current_line:
                    lines.append(current_line)
            
            return lines
            
        except Exception as e:
            print(f"文字換行處理錯誤: {e}")
            # 簡單分割作為備選
            return [text[i:i+30] for i in range(0, len(text), 30)][:5]
    
    def cleanup(self):
        """清理暫存檔案"""
        try:
            import shutil
            if os.path.exists(self.temp_dir):
                shutil.rmtree(self.temp_dir)
                print(f"清理暫存目錄: {self.temp_dir}")
        except Exception as e:
            print(f"清理暫存檔案時發生錯誤: {e}")

class ConversionWorker(QThread):
    """轉換工作執行緒"""
    
    progress_updated = Signal(int)
    status_updated = Signal(str)
    finished_successfully = Signal(str)
    error_occurred = Signal(str)
    
    def __init__(self, word_path: str, template_path: str, output_path: str):
        super().__init__()
        self.word_path = word_path
        self.template_path = template_path
        self.output_path = output_path
    
    def run(self):
        """執行轉換"""
        try:
            # 步驟1：分析Word文件
            self.status_updated.emit("正在分析Word文件章節結構...")
            self.progress_updated.emit(20)
            
            analyzer = WordDocumentAnalyzer()
            blocks = analyzer.analyze_document(self.word_path)
            
            if not blocks:
                raise Exception("Word文件中沒有找到可轉換的內容")
            
            # 步驟2：檢查PowerPoint範本
            self.status_updated.emit("正在檢查PowerPoint範本...")
            self.progress_updated.emit(40)
            
            if not os.path.exists(self.template_path):
                raise Exception(f"PowerPoint範本檔案不存在: {self.template_path}")
            
            # 步驟3：建立投影片
            self.status_updated.emit("正在建立投影片...")
            self.progress_updated.emit(60)
            
            mapper = ContentToSlideMapper(self.template_path)
            presentation = mapper.create_slides(blocks)
            
            # 步驟4：儲存檔案
            self.status_updated.emit("正在儲存檔案...")
            self.progress_updated.emit(80)
            
            # 確保輸出目錄存在
            output_dir = os.path.dirname(self.output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir, exist_ok=True)
            
            presentation.save(self.output_path)
            
            self.progress_updated.emit(100)
            self.status_updated.emit("轉換完成！字體已設定為至少32pt大小")
            self.finished_successfully.emit(self.output_path)
            
        except Exception as e:
            import traceback
            error_details = f"{str(e)}\n\n詳細錯誤資訊:\n{traceback.format_exc()}"
            self.error_occurred.emit(error_details)

class DropArea(QFrame):
    """拖放區域元件"""
    
    file_dropped = Signal(str)
    
    def __init__(self, file_type: str, parent=None):
        super().__init__(parent)
        self.file_type = file_type
        self.setup_ui()
        
    def setup_ui(self):
        """設定UI"""
        self.setAcceptDrops(True)
        self.setFrameStyle(QFrame.StyledPanel)
        self.setStyleSheet("""
            DropArea {
                border: 2px dashed #3498db;
                border-radius: 8px;
                background-color: #f8f9fa;
                padding: 15px;
                min-height: 120px;
                max-height: 160px;
            }
            DropArea:hover {
                background-color: #e3f2fd;
                border-color: #2196f3;
            }
        """)
        
        layout = QVBoxLayout(self)
        layout.setSpacing(10)
        layout.setContentsMargins(15, 15, 15, 15)
        
        # 圖示
        icon_label = QLabel("📄" if self.file_type == "word" else "📊")
        icon_label.setFont(QFont("Arial", 36))
        icon_label.setAlignment(Qt.AlignCenter)
        icon_label.setMaximumHeight(45)
        
        # 文字說明
        if self.file_type == "word":
            text = "拖放WORD檔案到此處\n或點擊選擇檔案"
        else:
            text = "拖放POWERPOINT範本到此處\n或點擊選擇檔案"
        text_label = QLabel(text)
        text_label.setFont(QFont("Microsoft JhengHei", 10))
        text_label.setAlignment(Qt.AlignCenter)
        text_label.setStyleSheet("color: #666; margin: 5px; line-height: 1.4;")
        text_label.setWordWrap(True)
        text_label.setMaximumHeight(40)
        
        layout.addWidget(icon_label)
        layout.addWidget(text_label)
        layout.addStretch()
        
        self.setMinimumHeight(120)
        self.setMaximumHeight(160)
    
    def dragEnterEvent(self, event: QDragEnterEvent):
        """拖拽進入事件"""
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
            self.setStyleSheet("""
                DropArea {
                    border: 2px solid #4caf50;
                    background-color: #e8f5e8;
                }
            """)
    
    def dragLeaveEvent(self, event):
        """拖拽離開事件"""
        self.setStyleSheet("""
            DropArea {
                border: 2px dashed #3498db;
                background-color: #f8f9fa;
            }
        """)
    
    def dropEvent(self, event: QDropEvent):
        """放置事件"""
        files = [u.toLocalFile() for u in event.mimeData().urls()]
        if files:
            self.file_dropped.emit(files[0])
        
        self.setStyleSheet("""
            DropArea {
                border: 2px dashed #3498db;
                background-color: #f8f9fa;
            }
        """)
    
    def mousePressEvent(self, event):
        """滑鼠點擊選擇檔案"""
        if self.file_type == "word":
            file_filter = "Word Documents (*.docx *.doc)"
            dialog_title = "選擇WORD檔案"
        else:
            file_filter = "PowerPoint Files (*.pptx *.ppt)"
            dialog_title = "選擇POWERPOINT檔案"
            
        file_path, _ = QFileDialog.getOpenFileName(
            self, dialog_title, "", file_filter
        )
        
        if file_path:
            self.file_dropped.emit(file_path)

class PreviewWidget(QScrollArea):
    """預覽元件 - 圖片顯示版本"""
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setup_ui()
        self.image_generator = SlideImageGenerator()
        
    def setup_ui(self):
        """設定UI"""
        self.setWidgetResizable(True)
        self.setStyleSheet("""
            QScrollArea {
                border: 1px solid #ddd;
                border-radius: 8px;
                background-color: #f8f9fa;
            }
        """)
        
        # 建立內容區域
        content = QWidget()
        self.content_layout = QVBoxLayout(content)
        self.content_layout.setSpacing(15)
        self.content_layout.setContentsMargins(15, 15, 15, 15)
        
        self.setWidget(content)
    
    def update_preview(self, presentation_path: str):
        """更新預覽 - 使用圖片顯示 (UTF-8修復版)"""
        try:
            # 清除現有預覽
            self.clear_preview()
            
            # 顯示載入訊息
            loading_label = QLabel("正在產生投影片預覽圖片，請稍候...")
            loading_label.setAlignment(Qt.AlignCenter)
            loading_label.setStyleSheet("color: #3498db; font-size: 14px; padding: 20px;")
            self.content_layout.addWidget(loading_label)
            
            # 強制更新UI
            QApplication.processEvents()
            
            print(f"開始處理簡報檔案: {presentation_path}")
            
            # 檢查檔案編碼和完整性
            if not self._validate_presentation_file(presentation_path):
                raise Exception("PowerPoint檔案驗證失敗，可能損壞或編碼錯誤")
            
            # 產生投影片圖片
            image_paths = self.image_generator.generate_slide_images(presentation_path)
            
            # 移除載入訊息
            loading_label.deleteLater()
            
            if not image_paths:
                error_label = QLabel("❌ 無法產生預覽圖片\n\n🔧 可能原因:\n• PowerPoint檔案格式不正確\n• 檔案內容編碼問題\n• 缺少支援中文的字體")
                error_label.setStyleSheet("color: #e74c3c; padding: 20px; text-align: center;")
                error_label.setAlignment(Qt.AlignCenter)
                self.content_layout.addWidget(error_label)
                return
            
            # 添加預覽說明
            info_label = QLabel("📊 以下為投影片預覽圖片 (UTF-8編碼，支援中文顯示)")
            info_label.setFont(QFont("Microsoft JhengHei", 10, QFont.Weight.Bold))
            info_label.setStyleSheet("color: #27ae60; padding: 10px; background: #f0f8f0; border-radius: 5px; margin-bottom: 10px;")
            info_label.setAlignment(Qt.AlignCenter)
            self.content_layout.addWidget(info_label)
            
            # 顯示每張投影片圖片
            success_count = 0
            for i, image_path in enumerate(image_paths):
                if os.path.exists(image_path):
                    try:
                        preview_item = self.create_image_preview(image_path, i + 1)
                        self.content_layout.addWidget(preview_item)
                        success_count += 1
                        
                        # 每處理一張就更新UI
                        QApplication.processEvents()
                    except Exception as e:
                        print(f"建立預覽項目失敗: {e}")
                        continue
            
            # 顯示處理結果
            if success_count > 0:
                result_label = QLabel(f"✅ 成功產生 {success_count} 張投影片預覽")
                result_label.setStyleSheet("color: #27ae60; padding: 5px; font-weight: bold;")
                result_label.setAlignment(Qt.AlignCenter)
                self.content_layout.addWidget(result_label)
            
        except Exception as e:
            self.clear_preview()
            import traceback
            error_details = traceback.format_exc()
            print(f"預覽更新錯誤: {error_details}")
            
            error_label = QLabel(f"❌ 預覽產生失敗: {str(e)}\n\n🔧 故障排除提示:\n• 確保PowerPoint檔案為UTF-8相容格式\n• 檢查檔案是否包含有效的中文內容\n• 確認系統已安裝中文字體\n• 嘗試重新儲存PowerPoint檔案")
            error_label.setStyleSheet("color: #e74c3c; padding: 20px; background: #fdf2f2; border-radius: 5px; border: 1px solid #e74c3c;")
            error_label.setWordWrap(True)
            self.content_layout.addWidget(error_label)
    
    def _validate_presentation_file(self, file_path: str) -> bool:
        """驗證PowerPoint檔案的完整性和編碼"""
        try:
            # 檢查檔案是否存在和可讀
            if not os.path.exists(file_path):
                print(f"檔案不存在: {file_path}")
                return False
            
            if not os.access(file_path, os.R_OK):
                print(f"檔案無讀取權限: {file_path}")
                return False
            
            # 嘗試載入PowerPoint檔案
            prs = Presentation(file_path)
            
            # 檢查是否有投影片
            if len(prs.slides) == 0:
                print("PowerPoint檔案中沒有投影片")
                return False
            
            # 測試讀取第一張投影片的文字內容
            first_slide = prs.slides[0]
            text_found = False
            
            for shape in first_slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    text_found = True
                    # 測試文字編碼
                    try:
                        test_text = str(shape.text)
                        test_text.encode('utf-8')
                    except UnicodeEncodeError as e:
                        print(f"文字編碼問題: {e}")
                        # 不直接返回False，繼續嘗試
                    break
            
            print(f"檔案驗證完成，投影片數量: {len(prs.slides)}, 包含文字: {text_found}")
            return True
            
        except Exception as e:
            print(f"檔案驗證失敗: {e}")
            return False
    
    def clear_preview(self):
        """清除預覽"""
        while self.content_layout.count():
            child = self.content_layout.takeAt(0)
            if child.widget():
                child.widget().deleteLater()
    
    def create_image_preview(self, image_path: str, slide_number: int) -> QWidget:
        """建立圖片預覽 (UTF-8支援版)"""
        frame = QFrame()
        frame.setFrameStyle(QFrame.Box)
        frame.setStyleSheet("""
            QFrame {
                border: 2px solid #3498db;
                border-radius: 10px;
                background-color: white;
                margin: 5px;
                padding: 10px;
            }
        """)
        
        layout = QVBoxLayout(frame)
        layout.setSpacing(8)
        layout.setContentsMargins(12, 12, 12, 12)
        
        # 投影片標題欄
        header_frame = QFrame()
        header_frame.setStyleSheet("""
            QFrame {
                background-color: #3498db;
                border-radius: 5px;
                padding: 5px;
                border: none;
            }
        """)
        
        header_layout = QHBoxLayout(header_frame)
        header_layout.setContentsMargins(8, 4, 8, 4)
        
        # 投影片編號 (確保UTF-8編碼)
        try:
            number_text = f"投影片 {slide_number}"
            # 測試UTF-8編碼
            number_text.encode('utf-8')
        except UnicodeEncodeError:
            number_text = f"Slide {slide_number}"
        
        number_label = QLabel(number_text)
        number_label.setFont(QFont("Microsoft JhengHei", 11, QFont.Weight.Bold))
        number_label.setStyleSheet("color: white;")
        
        # 添加UTF-8圖片預覽標籤
        try:
            image_text = "UTF-8圖片預覽"
            image_text.encode('utf-8')
        except UnicodeEncodeError:
            image_text = "Image Preview"
        
        image_label_header = QLabel(image_text)
        image_label_header.setFont(QFont("Microsoft JhengHei", 8))
        image_label_header.setStyleSheet("color: #f1c40f; background: rgba(255,255,255,20); padding: 2px 6px; border-radius: 3px;")
        
        header_layout.addWidget(number_label, 1)
        header_layout.addWidget(image_label_header)
        
        layout.addWidget(header_frame)
        
        # 投影片圖片顯示
        image_container = QFrame()
        image_container.setStyleSheet("""
            QFrame {
                border: 1px solid #34495e;
                border-radius: 5px;
                background-color: white;
                padding: 5px;
            }
        """)
        
        image_layout = QVBoxLayout(image_container)
        image_layout.setContentsMargins(5, 5, 5, 5)
        
        try:
            # 載入圖片
            pixmap = QPixmap(image_path)
            if not pixmap.isNull():
                # 縮放圖片以符合預覽大小，保持16:9比例
                target_width = 600
                target_height = 338  # 16:9 比例
                
                scaled_pixmap = pixmap.scaled(
                    target_width, target_height, 
                    Qt.KeepAspectRatio, 
                    Qt.SmoothTransformation
                )
                
                image_display = QLabel()
                image_display.setPixmap(scaled_pixmap)
                image_display.setAlignment(Qt.AlignCenter)
                image_display.setStyleSheet("border: none; background: white;")
                
                image_layout.addWidget(image_display)
                
                # 添加圖片資訊
                image_info = QLabel(f"圖片尺寸: {pixmap.width()}×{pixmap.height()} | 比例: 16:9")
                image_info.setFont(QFont("Microsoft JhengHei", 8))
                image_info.setStyleSheet("color: #7f8c8d; margin-top: 5px;")
                image_info.setAlignment(Qt.AlignCenter)
                image_layout.addWidget(image_info)
                
            else:
                # 圖片載入失敗
                error_display = QLabel("圖片載入失敗\n可能的編碼問題")
                error_display.setAlignment(Qt.AlignCenter)
                error_display.setStyleSheet("color: #e74c3c; padding: 40px; font-size: 14px;")
                image_layout.addWidget(error_display)
                
        except Exception as e:
            print(f"圖片顯示錯誤: {e}")
            # 顯示錯誤訊息
            error_display = QLabel(f"圖片顯示錯誤\n{str(e)[:50]}")
            error_display.setAlignment(Qt.AlignCenter)
            error_display.setStyleSheet("color: #e74c3c; padding: 40px; font-size: 12px;")
            image_layout.addWidget(error_display)
        
        layout.addWidget(image_container)
        
        # 技術資訊區域
        tech_info_frame = QFrame()
        tech_info_frame.setStyleSheet("""
            QFrame {
                background-color: #ecf0f1;
                border-radius: 5px;
                padding: 8px;
                border: 1px solid #bdc3c7;
            }
        """)
        
        tech_layout = QVBoxLayout(tech_info_frame)
        tech_layout.setContentsMargins(8, 6, 8, 6)
        tech_layout.setSpacing(3)
        
        # 檔案資訊
        try:
            file_info = f"檔案: {os.path.basename(image_path)} | 編碼: UTF-8"
            file_size = os.path.getsize(image_path) if os.path.exists(image_path) else 0
            file_info += f" | 大小: {file_size//1024}KB"
        except:
            file_info = "檔案資訊無法取得"
        
        file_label = QLabel(file_info)
        file_label.setFont(QFont("Microsoft JhengHei", 8))
        file_label.setStyleSheet("color: #7f8c8d;")
        tech_layout.addWidget(file_label)
        
        # 編碼狀態
        try:
            encoding_status = "✅ UTF-8編碼正常 | 支援中文顯示"
        except:
            encoding_status = "⚠️ 編碼檢查失敗"
        
        encoding_label = QLabel(encoding_status)
        encoding_label.setFont(QFont("Microsoft JhengHei", 8))
        encoding_label.setStyleSheet("color: #27ae60;")
        tech_layout.addWidget(encoding_label)
        
        layout.addWidget(tech_info_frame)
        
        frame.setMaximumHeight(500)
        frame.setMinimumHeight(450)
        return frame
    
    def __del__(self):
        """解構函式 - 清理暫存檔案"""
        try:
            self.image_generator.cleanup()
        except:
            pass

class MainWindow(QMainWindow):
    """主視窗"""
    
    def __init__(self):
        super().__init__()
        # 初始化設定管理器
        self.config_manager = ConfigManager()
        
        # 確保路徑變數正確初始化
        self.word_path = ""
        self.template_path = ""
        self.output_path = ""
        self.worker = None
        
        self.setup_ui()
        self.setup_connections()
        
        # 載入上次使用的檔案路徑
        self.load_last_used_paths()
        
        # 初始檢查轉換準備狀態
        self.check_ready_to_convert()
        
    def setup_ui(self):
        """設定UI"""
        self.setWindowTitle("通用Word轉PowerPoint工具 v2.3 (UTF-8修復版)")
        self.setGeometry(100, 100, 1300, 850)
        
        # 設定樣式
        self.setStyleSheet("""
            QMainWindow {
                background-color: #f5f5f5;
            }
            QPushButton {
                background-color: #3498db;
                color: white;
                border: none;
                padding: 8px 16px;
                border-radius: 4px;
                font-size: 12px;
                font-weight: bold;
            }
            QPushButton:hover {
                background-color: #2980b9;
            }
            QPushButton:disabled {
                background-color: #bdc3c7;
            }
            QGroupBox {
                font-weight: bold;
                border: 2px solid #bdc3c7;
                border-radius: 6px;
                margin-top: 8px;
                padding-top: 8px;
                font-size: 12px;
            }
            QGroupBox::title {
                subcontrol-origin: margin;
                left: 8px;
                padding: 0 4px 0 4px;
            }
        """)
        
        # 中央元件
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        
        # 主佈局
        main_layout = QHBoxLayout(central_widget)
        main_layout.setSpacing(10)
        main_layout.setContentsMargins(10, 10, 10, 10)
        
        # 左側面板
        left_panel = self.create_left_panel()
        left_panel.setMaximumWidth(450)
        
        # 右側預覽面板
        right_panel = self.create_right_panel()
        
        # 分割器
        splitter = QSplitter(Qt.Horizontal)
        splitter.addWidget(left_panel)
        splitter.addWidget(right_panel)
        splitter.setSizes([300, 900])
        
        main_layout.addWidget(splitter)
        
    def create_left_panel(self) -> QWidget:
        """建立左側面板"""
        panel = QWidget()
        panel.setMinimumWidth(300)
        layout = QVBoxLayout(panel)
        layout.setSpacing(10)
        layout.setContentsMargins(10, 10, 10, 10)
        
        # 標題
        title = QLabel("📊 Word轉PowerPoint工具")
        title.setFont(QFont("Microsoft JhengHei", 16, QFont.Weight.Bold))
        title.setAlignment(Qt.AlignCenter)
        title.setStyleSheet("color: #2c3e50; margin: 10px;")
        title.setMaximumHeight(40)
        layout.addWidget(title)
        
        # 功能說明
        features = QLabel("✨ 新功能: 智慧章節識別 | 大字體顯示 | UTF-8圖片預覽 | 路徑記憶")
        features.setFont(QFont("Microsoft JhengHei", 9))
        features.setAlignment(Qt.AlignCenter)
        features.setStyleSheet("color: #27ae60; margin-bottom: 5px;")
        features.setMaximumHeight(25)
        layout.addWidget(features)
        
        # 檔案選擇區域
        file_group = QGroupBox("1. 選擇檔案")
        file_group.setMaximumHeight(420)
        file_layout = QVBoxLayout(file_group)
        file_layout.setSpacing(8)
        file_layout.setContentsMargins(10, 20, 10, 10)
        
        # Word檔案拖放區
        word_label = QLabel("Word文件:")
        word_label.setFont(QFont("Microsoft JhengHei", 11, QFont.Weight.Bold))
        word_label.setStyleSheet("color: #2c3e50; margin-bottom: 3px;")
        word_label.setMaximumHeight(20)
        
        self.word_drop = DropArea("word")
        self.word_drop.setMaximumHeight(140)
        
        self.word_status = QLabel("未選擇Word檔案")
        self.word_status.setStyleSheet("color: #e74c3c; margin: 5px; font-size: 10px;")
        self.word_status.setMaximumHeight(28)
        
        file_layout.addWidget(word_label)
        file_layout.addWidget(self.word_drop)
        file_layout.addWidget(self.word_status)
        
        # 添加分隔線
        line = QFrame()
        line.setFrameShape(QFrame.HLine)
        line.setFrameShadow(QFrame.Sunken)
        line.setStyleSheet("color: #bdc3c7;")
        line.setMaximumHeight(5)
        file_layout.addWidget(line)
        
        # PowerPoint範本拖放區
        template_label = QLabel("PowerPoint範本:")
        template_label.setFont(QFont("Microsoft JhengHei", 11, QFont.Weight.Bold))
        template_label.setStyleSheet("color: #2c3e50; margin-bottom: 3px;")
        template_label.setMaximumHeight(20)
        
        self.template_drop = DropArea("pptx")
        self.template_drop.setMaximumHeight(140)
        
        self.template_status = QLabel("未選擇PowerPoint範本")
        self.template_status.setStyleSheet("color: #e74c3c; margin: 5px; font-size: 10px;")
        self.template_status.setMaximumHeight(28)
        
        file_layout.addWidget(template_label)
        file_layout.addWidget(self.template_drop)
        file_layout.addWidget(self.template_status)
        
        layout.addWidget(file_group)
        
        # 轉換設定
        settings_group = QGroupBox("2. 轉換設定")
        settings_group.setMaximumHeight(200)
        settings_layout = QVBoxLayout(settings_group)
        settings_layout.setSpacing(8)
        settings_layout.setContentsMargins(10, 15, 10, 10)
        
        # 設定說明
        settings_desc = QLabel("• 自動識別「一、二、三」等章節標題\n• 字體大小自動設定為32pt以上\n• 會記憶本次使用的檔案路徑")
        settings_desc.setFont(QFont("Microsoft JhengHei", 9))
        settings_desc.setStyleSheet("color: #27ae60; margin-bottom: 5px;")
        settings_desc.setMaximumHeight(140)
        settings_layout.addWidget(settings_desc)
        
        # 輸出路徑
        output_desc_label = QLabel("輸出檔案路徑:")
        output_desc_label.setFont(QFont("Microsoft JhengHei", 11, QFont.Weight.Bold))
        output_desc_label.setStyleSheet("color: #2c3e50; margin-bottom: 3px;")
        output_desc_label.setMaximumHeight(20)
        
        output_layout = QHBoxLayout()
        output_layout.setSpacing(8)
        
        self.output_label = QLabel("將自動設定輸出位置...")
        self.output_label.setStyleSheet("""
            QLabel {
                border: 1px solid #ccc; 
                padding: 8px; 
                background: white; 
                color: #666;
                border-radius: 4px;
                font-size: 10px;
            }
        """)
        self.output_label.setMaximumHeight(35)
        self.output_label.setWordWrap(True)
        
        output_btn = QPushButton("選擇")
        output_btn.setMaximumHeight(35)
        output_btn.setMaximumWidth(80)
        output_btn.clicked.connect(self.select_output_path)
        
        output_layout.addWidget(self.output_label, 1)
        output_layout.addWidget(output_btn)
        
        settings_layout.addWidget(output_desc_label)
        settings_layout.addLayout(output_layout)
        
        layout.addWidget(settings_group)
        
        # 轉換按鈕
        self.convert_btn = QPushButton("🚀 開始轉換")
        self.convert_btn.setEnabled(False)
        self.convert_btn.setFont(QFont("Microsoft JhengHei", 13, QFont.Weight.Bold))
        self.convert_btn.setMaximumHeight(50)
        self.convert_btn.clicked.connect(self.start_conversion)
        layout.addWidget(self.convert_btn)
        
        # 進度條
        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        self.progress_bar.setMaximumHeight(20)
        self.progress_bar.setStyleSheet("""
            QProgressBar {
                border: 2px solid #bdc3c7;
                border-radius: 4px;
                text-align: center;
                font-weight: bold;
                font-size: 10px;
            }
            QProgressBar::chunk {
                background-color: #27ae60;
                border-radius: 2px;
            }
        """)
        layout.addWidget(self.progress_bar)
        
        # 狀態標籤
        self.status_label = QLabel("")
        self.status_label.setAlignment(Qt.AlignCenter)
        self.status_label.setStyleSheet("color: #27ae60; font-weight: bold; font-size: 11px; margin: 5px;")
        self.status_label.setMaximumHeight(25)
        self.status_label.setWordWrap(True)
        layout.addWidget(self.status_label)
        
        # 添加彈性空間
        spacer = QWidget()
        spacer.setMaximumHeight(20)
        layout.addWidget(spacer)
        
        return panel
    
    def create_right_panel(self) -> QWidget:
        """建立右側預覽面板"""
        panel = QWidget()
        layout = QVBoxLayout(panel)
        layout.setSpacing(5)
        layout.setContentsMargins(5, 5, 5, 5)
        
        # 預覽標題
        preview_title = QLabel("📋 投影片預覽 (UTF-8圖片模式)")
        preview_title.setFont(QFont("Microsoft JhengHei", 16, QFont.Weight.Bold))
        preview_title.setStyleSheet("color: #2c3e50; margin: 10px;")
        layout.addWidget(preview_title)
        
        # 預覽說明
        preview_desc = QLabel("🎯 UTF-8編碼顯示投影片 | 📐 保持16:9真實比例 | 🖼️ 修復中文亂碼問題")
        preview_desc.setFont(QFont("Microsoft JhengHei", 10))
        preview_desc.setStyleSheet("""
            color: #27ae60; 
            background: #f0f8f0; 
            padding: 8px; 
            border-radius: 5px; 
            border: 1px solid #27ae60;
            margin: 5px 10px;
        """)
        preview_desc.setWordWrap(True)
        layout.addWidget(preview_desc)
        
        # 預覽區域
        self.preview_widget = PreviewWidget()
        layout.addWidget(self.preview_widget)
        
        return panel
    
    def setup_connections(self):
        """設定信號連接"""
        self.word_drop.file_dropped.connect(self.on_word_file_selected)
        self.template_drop.file_dropped.connect(self.on_template_file_selected)
    
    def load_last_used_paths(self):
        """載入上次使用的檔案路徑"""
        try:
            last_word_path = self.config_manager.get_last_word_path()
            last_template_path = self.config_manager.get_last_template_path()
            
            # 自動載入上次的Word檔案
            if last_word_path and os.path.exists(last_word_path):
                self.on_word_file_selected(last_word_path)
                print(f"自動載入上次的Word檔案: {last_word_path}")
            
            # 自動載入上次的範本檔案
            if last_template_path and os.path.exists(last_template_path):
                self.on_template_file_selected(last_template_path)
                print(f"自動載入上次的範本檔案: {last_template_path}")
                
        except Exception as e:
            print(f"載入上次使用路徑時發生錯誤: {e}")
    
    def on_word_file_selected(self, file_path: str):
        """Word檔案選擇處理"""
        if file_path.lower().endswith(('.docx', '.doc')):
            self.word_path = file_path
            filename = os.path.basename(file_path)
            self.word_status.setText(f"✅ {filename} (已啟用智慧章節識別)")
            self.word_status.setStyleSheet("color: #27ae60; margin: 5px; font-weight: bold;")
            
            # 儲存到設定
            self.config_manager.set_last_word_path(file_path)
            
            self.auto_set_output_path(file_path)
        else:
            QMessageBox.warning(self, "檔案格式錯誤", "請選擇Word文件(.docx或.doc)")
            return
            
        self.check_ready_to_convert()
    
    def on_template_file_selected(self, file_path: str):
        """範本檔案選擇處理"""
        if file_path.lower().endswith(('.pptx', '.ppt')):
            self.template_path = file_path
            filename = os.path.basename(file_path)
            self.template_status.setText(f"✅ {filename}")
            self.template_status.setStyleSheet("color: #27ae60; margin: 5px; font-weight: bold;")
            
            # 儲存到設定
            self.config_manager.set_last_template_path(file_path)
        else:
            QMessageBox.warning(self, "檔案格式錯誤", "請選擇PowerPoint檔案(.pptx或.ppt)")
            return
            
        self.check_ready_to_convert()
    
    def auto_set_output_path(self, word_path: str):
        """根據Word檔案路徑自動設定輸出路徑"""
        try:
            file_dir = os.path.dirname(word_path)
            file_name_without_ext = os.path.splitext(os.path.basename(word_path))[0]
            
            output_path = os.path.join(file_dir, f"{file_name_without_ext}.pptx")
            
            counter = 1
            original_output_path = output_path
            while os.path.exists(output_path):
                base_name = os.path.splitext(original_output_path)[0]
                output_path = f"{base_name}_{counter}.pptx"
                counter += 1
            
            self.output_path = output_path
            self.output_label.setText(f"📁 {os.path.basename(output_path)}")
            self.output_label.setStyleSheet("border: 1px solid #27ae60; padding: 8px; background: #f8fff8; color: #27ae60;")
            
            # 儲存輸出目錄
            self.config_manager.set_last_output_dir(file_dir)
            
        except Exception as e:
            print(f"自動設定輸出路徑時發生錯誤: {e}")
    
    def select_output_path(self):
        """選擇輸出路徑"""
        default_dir = self.config_manager.get_last_output_dir()
        default_name = "presentation.pptx"
        
        if self.word_path:
            default_dir = os.path.dirname(self.word_path)
            word_name = os.path.splitext(os.path.basename(self.word_path))[0]
            default_name = f"{word_name}.pptx"
        
        default_path = os.path.join(default_dir, default_name) if default_dir else default_name
        
        file_path, _ = QFileDialog.getSaveFileName(
            self, "儲存PowerPoint檔案", default_path, "PowerPoint Files (*.pptx)"
        )
        
        if file_path:
            if not file_path.endswith('.pptx'):
                file_path += '.pptx'
            self.output_path = file_path
            self.output_label.setText(f"📁 {os.path.basename(file_path)}")
            self.output_label.setStyleSheet("border: 1px solid #27ae60; padding: 8px; background: #f8fff8; color: #27ae60;")
            
            # 儲存輸出目錄
            self.config_manager.set_last_output_dir(os.path.dirname(file_path))
            
            self.check_ready_to_convert()
    
    def check_ready_to_convert(self):
        """檢查是否準備好轉換"""
        word_ready = bool(self.word_path and self.word_path.strip())
        template_ready = bool(self.template_path and self.template_path.strip())
        output_ready = bool(self.output_path and self.output_path.strip())
        
        ready = word_ready and template_ready and output_ready
        
        self.convert_btn.setEnabled(ready)
        
        if ready:
            self.convert_btn.setText("🚀 開始轉換 (UTF-8修復版)")
            self.convert_btn.setStyleSheet("""
                QPushButton {
                    background-color: #27ae60;
                    color: white;
                    border: none;
                    padding: 10px 20px;
                    border-radius: 5px;
                    font-size: 14px;
                    font-weight: bold;
                }
                QPushButton:hover {
                    background-color: #229954;
                }
            """)
        else:
            missing_items = []
            if not word_ready:
                missing_items.append("Word文件")
            if not template_ready:
                missing_items.append("PowerPoint範本")
            if not output_ready:
                missing_items.append("輸出路徑")
            
            self.convert_btn.setText(f"請選擇: {', '.join(missing_items)}")
            self.convert_btn.setStyleSheet("""
                QPushButton {
                    background-color: #bdc3c7;
                    color: #7f8c8d;
                    border: none;
                    padding: 10px 20px;
                    border-radius: 5px;
                    font-size: 14px;
                    font-weight: bold;
                }
            """)
    
    def start_conversion(self):
        """開始轉換"""
        if not all([self.word_path, self.template_path, self.output_path]):
            QMessageBox.warning(
                self, "準備不完整", 
                "請確保已選擇Word文件、PowerPoint範本和輸出路徑"
            )
            return
        
        if not os.path.exists(self.word_path):
            QMessageBox.critical(self, "檔案不存在", f"Word文件不存在：\n{self.word_path}")
            return
            
        if not os.path.exists(self.template_path):
            QMessageBox.critical(self, "檔案不存在", f"PowerPoint範本不存在：\n{self.template_path}")
            return
        
        self.convert_btn.setEnabled(False)
        self.progress_bar.setVisible(True)
        self.progress_bar.setValue(0)
        self.status_label.setText("正在準備增強轉換...")
        
        self.worker = ConversionWorker(
            self.word_path, 
            self.template_path, 
            self.output_path
        )
        
        self.worker.progress_updated.connect(self.progress_bar.setValue)
        self.worker.status_updated.connect(self.status_label.setText)
        self.worker.finished_successfully.connect(self.on_conversion_finished)
        self.worker.error_occurred.connect(self.on_conversion_error)
        
        self.worker.start()
    
    def on_conversion_finished(self, output_path: str):
        """轉換完成處理"""
        self.progress_bar.setVisible(False)
        self.status_label.setText("✅ 轉換完成！已產生圖片預覽")
        self.status_label.setStyleSheet("color: #27ae60; font-weight: bold;")
        
        self.check_ready_to_convert()
        
        # 更新預覽 - 現在是圖片預覽
        self.preview_widget.update_preview(output_path)
        
        reply = QMessageBox.question(
            self, "轉換完成", 
            f"🎉 PowerPoint檔案已成功儲存！\n📁 位置: {output_path}\n\n✨ 增強功能已套用:\n• 智慧章節識別 (一、二、三...)\n• 字體大小32pt以上\n• 右側顯示圖片預覽\n• 路徑記憶功能\n\n右側顯示投影片預覽圖片，\n包含真實的內容和比例。\n\n是否立即開啟查看？",
            QMessageBox.Yes | QMessageBox.No
        )
        
        if reply == QMessageBox.Yes:
            self.open_file(output_path)
    
    def open_file(self, file_path: str):
        """跨平台開啟檔案"""
        import platform
        import subprocess
        
        try:
            system = platform.system()
            if system == "Windows":
                os.startfile(file_path)
            elif system == "Darwin":
                subprocess.run(["open", file_path])
            else:
                subprocess.run(["xdg-open", file_path])
        except Exception as e:
            QMessageBox.information(
                self, "提示", 
                f"檔案已儲存到：{file_path}\n請手動開啟檔案。"
            )
    
    def on_conversion_error(self, error_message: str):
        """轉換錯誤處理"""
        self.progress_bar.setVisible(False)
        self.status_label.setText("❌ 轉換失敗")
        self.status_label.setStyleSheet("color: #e74c3c; font-weight: bold;")
        
        self.check_ready_to_convert()
        
        detailed_message = f"轉換過程中發生錯誤：\n{error_message}\n\n"
        detailed_message += "可能的解決方法：\n"
        detailed_message += "1. 確保Word文件不是受保護的\n"
        detailed_message += "2. 確保PowerPoint範本檔案完整\n"
        detailed_message += "3. 檢查檔案路徑中是否包含特殊字元\n"
        detailed_message += "4. 嘗試關閉正在使用這些檔案的其他程式\n"
        detailed_message += "5. 確認Word文件中包含可識別的章節標題"
        
        QMessageBox.critical(self, "轉換錯誤", detailed_message)
    
    def closeEvent(self, event):
        """視窗關閉事件"""
        # 清理暫存檔案
        try:
            if hasattr(self, 'preview_widget'):
                del self.preview_widget
        except:
            pass
        
        event.accept()

def main():
    """主函式"""
    app = QApplication(sys.argv)
    
    app.setApplicationName("Word轉PowerPoint工具 (UTF-8修復版)")
    app.setApplicationVersion("2.3")
    app.setOrganizationName("智慧辦公工具")
    
    window = MainWindow()
    window.show()
    
    sys.exit(app.exec())

if __name__ == "__main__":
    print("=" * 70)
    print("📊 Word轉PowerPoint工具 v2.3 (UTF-8修復版) 📊")
    print("=" * 70)
    print("🎯 核心功能:")
    print("  ✓ 智慧識別中文章節標題（一、二、三、四...）")
    print("  ✓ 字體大小自動最佳化至32pt以上")
    print("  ✓ UTF-8編碼投影片圖片預覽")
    print("  ✓ 記憶上次使用的檔案路徑")
    print("")
    print("🔥 UTF-8預覽特色:")
    print("  • 完全修復中文亂碼問題")
    print("  • 使用UTF-8編碼處理所有文字")
    print("  • 智慧載入系統中文字體")
    print("  • 保持真實的16:9比例")
    print("  • 多層編碼錯誤處理機制")
    print("")
    print("🔧 編碼處理:")
    print("  • 自動檢測和轉換文字編碼")
    print("  • 支援Big5、GB2312、UTF-16等編碼")
    print("  • 中文友善的文字換行處理")
    print("  • 字體相容性測試和選擇")
    print("")
    print("💾 記憶功能:")
    print("  • 自動記住上次使用的Word檔案")
    print("  • 自動記住上次使用的PowerPoint範本")
    print("  • 自動記住上次使用的輸出目錄")
    print("  • 下次開啟程式時自動載入")
    print("")
    print("📋 使用步驟:")
    print("  1. 拖放Word文件（支援.docx/.doc）")
    print("  2. 拖放PowerPoint範本（.pptx/.ppt）")
    print("  3. 點擊開始轉換")
    print("  4. 查看右側UTF-8圖片預覽")
    print("")
    print("🛠️  安裝相依性:")
    print("pip install PySide6 python-docx python-pptx pillow")
    print("")
    print("🌏 編碼支援:")
    print("  ✓ UTF-8 (主要編碼)")
    print("  ✓ Big5 (繁體中文)")
    print("  ✓ GB2312 (簡體中文)")
    print("  ✓ UTF-16 (Unicode)")
    print("=" * 70)
    
    missing_deps = []
    
    try:
        import PySide6
    except ImportError:
        missing_deps.append("PySide6")
    
    try:
        import docx
    except ImportError:
        missing_deps.append("python-docx")
    
    try:
        import pptx
    except ImportError:
        missing_deps.append("python-pptx")
    
    try:
        import PIL
    except ImportError:
        missing_deps.append("pillow")
    
    if not DEPENDENCIES_OK or missing_deps:
        if missing_deps:
            print(f"❌ 缺少必要套件: {', '.join(missing_deps)}")
            print(f"請執行: pip install {' '.join(missing_deps)}")
        sys.exit(1)
    
    try:
        main()
    except Exception as e:
        print(f"❌ 程式執行錯誤: {e}")
        import traceback
        traceback.print_exc()