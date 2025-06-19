#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Word轉PowerPoint核心模組 - 支援 Streamlit 和獨立應用程式
核心功能模組，不依賴特定GUI框架
"""

import os
import re
import json
import platform
import tempfile
from pathlib import Path
from typing import List, Dict, Tuple, Optional
from dataclasses import dataclass
from io import BytesIO

# 第三方庫導入
try:
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance
    DEPENDENCIES_OK = True
except ImportError as e:
    print(f"請安裝必要套件: pip install python-docx python-pptx pillow")
    DEPENDENCIES_OK = False

@dataclass
class ContentBlock:
    """內容塊資料結構"""
    text: str
    level: int  # 標題層級 (0=主標題/章節, 1=次標題, 2=內容)
    content_type: str  # header, chapter, title, subtitle, content, quote, list
    formatting: Dict = None
    estimated_length: int = 0  # 估算文字長度

@dataclass
class ConversionResult:
    """轉換結果資料結構"""
    success: bool
    output_path: str = ""
    preview_images: List[str] = None
    error_message: str = ""
    slides_count: int = 0

class SystemFontManager:
    """跨平台系統字體管理器"""
    
    def __init__(self):
        self.platform = platform.system()
        self.available_fonts = {}
        self.font_cache = {}
        self._load_system_fonts()
    
    def _load_system_fonts(self):
        """載入系統字體並測試中文支援"""
        if self.platform == "Darwin":  # macOS
            font_candidates = [
                ("PingFang TC", [
                    "/System/Library/Fonts/PingFang.ttc",
                    "/System/Library/Fonts/PingFangTC-Regular.otf",
                    "/Library/Fonts/PingFang.ttc"
                ]),
                ("Hiragino Sans", [
                    "/System/Library/Fonts/Hiragino Sans GB.ttc",
                    "/Library/Fonts/Hiragino Sans GB.ttc"
                ]),
                ("STHeiti", [
                    "/System/Library/Fonts/STHeiti Medium.ttc",
                    "/Library/Fonts/STHeiti Medium.ttc"
                ]),
                ("Arial Unicode MS", [
                    "/Library/Fonts/Arial Unicode MS.ttf",
                    "/System/Library/Fonts/Arial Unicode MS.ttf"
                ])
            ]
        elif self.platform == "Windows":  # Windows
            font_candidates = [
                ("Microsoft JhengHei", [
                    "C:/Windows/Fonts/msjh.ttc",
                    "C:/Windows/Fonts/msjhbd.ttc"
                ]),
                ("Microsoft YaHei", [
                    "C:/Windows/Fonts/msyh.ttc", 
                    "C:/Windows/Fonts/msyhbd.ttc"
                ]),
                ("SimSun", [
                    "C:/Windows/Fonts/simsun.ttc"
                ]),
                ("Arial Unicode MS", [
                    "C:/Windows/Fonts/ARIALUNI.TTF"
                ])
            ]
        else:  # Linux
            font_candidates = [
                ("WenQuanYi", [
                    "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
                    "/usr/share/fonts/truetype/arphic/uming.ttc"
                ]),
                ("Noto CJK", [
                    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
                    "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc"
                ])
            ]
        
        # 測試每種字體是否支援中文
        for font_name, paths in font_candidates:
            for path in paths:
                if os.path.exists(path):
                    if self._test_font_chinese_support(path):
                        self.available_fonts[font_name] = path
                        print(f"載入中文字體: {font_name} ({path})")
                        break
        
        if not self.available_fonts:
            print(f"警告: 在 {self.platform} 系統上未找到支援中文的字體")
    
    def _test_font_chinese_support(self, font_path: str) -> bool:
        """測試字體是否支援中文"""
        try:
            test_font = ImageFont.truetype(font_path, 20)
            test_img = Image.new('RGB', (100, 50), 'white')
            test_draw = ImageDraw.Draw(test_img)
            test_text = "測試中文字體"
            test_draw.text((10, 10), test_text, font=test_font, fill='black')
            return True
        except Exception as e:
            return False
    
    def get_best_font(self, size: int) -> ImageFont.ImageFont:
        """取得最佳的中文字體"""
        cache_key = f"font_{size}_{self.platform}"
        if cache_key in self.font_cache:
            return self.font_cache[cache_key]
        
        # 依平台優先順序嘗試字體
        if self.platform == "Darwin":
            priority_fonts = ["PingFang TC", "Hiragino Sans", "STHeiti", "Arial Unicode MS"]
        elif self.platform == "Windows":
            priority_fonts = ["Microsoft JhengHei", "Microsoft YaHei", "SimSun", "Arial Unicode MS"]
        else:
            priority_fonts = ["WenQuanYi", "Noto CJK"]
        
        for font_name in priority_fonts:
            if font_name in self.available_fonts:
                try:
                    font = ImageFont.truetype(self.available_fonts[font_name], size)
                    self.font_cache[cache_key] = font
                    return font
                except Exception as e:
                    print(f"載入字體失敗 {font_name}: {e}")
                    continue
        
        # 如果都失敗，使用預設字體
        default_font = ImageFont.load_default()
        self.font_cache[cache_key] = default_font
        return default_font

class WordDocumentAnalyzer:
    """Word文件智慧分析器"""
    
    def __init__(self):
        # 增強的中文章節識別模式
        self.chapter_patterns = [
            r'^[一二三四五六七八九十]+[、．.]\s*',
            r'^第[一二三四五六七八九十壹貳參肆伍陸柒捌玖拾]+[章節部分]\s*',
            r'^第[一二三四五六七八九十]+[、．.]\s*',
            r'^前言\s*',
            r'^結論\s*',
            r'^總結\s*',
            r'^概述\s*',
            r'^摘要\s*',
            r'^序言\s*',
            r'^引言\s*',
        ]
        
        self.subtitle_patterns = [
            r'^[一二三四五六七八九十]+[）)]\s*',
            r'^\([一二三四五六七八九十]+\)\s*',
            r'^[1-9]\d*[）)]\s*',
            r'^\([1-9]\d*\)\s*',
            r'^[a-z][）)]\s*',
            r'^\([a-z]\)\s*',
            r'^[•·○]\s*',
        ]
        
    def analyze_document(self, file_path: str = None, file_content: bytes = None) -> List[ContentBlock]:
        """分析Word文件結構，支援檔案路徑或檔案內容"""
        try:
            if file_content:
                # 從bytes內容載入文件（用於Streamlit上傳的檔案）
                doc = Document(BytesIO(file_content))
            elif file_path:
                # 從檔案路徑載入文件
                doc = Document(file_path)
            else:
                raise ValueError("必須提供 file_path 或 file_content")
            
            blocks = []
            
            header_line = True
            for para in doc.paragraphs:
                if not para.text.strip():
                    continue

                if header_line:
                    block = self._classify_header(para)
                    header_line = False
                else:
                    block = self._analyze_paragraph(para)
                if block:
                    # 估算文字長度（用於後續分頁判斷）
                    block.estimated_length = self._estimate_text_length(block.text)
                    blocks.append(block)
            
            return self._optimize_structure(blocks)
            
        except Exception as e:
            raise Exception(f"分析Word文件時發生錯誤: {e}")
    
    def _estimate_text_length(self, text: str) -> int:
        """估算文字在投影片上的顯示長度"""
        # 中文字符算2個單位，英文算1個單位
        length = 0
        for char in text:
            if ord(char) > 127:  # 非ASCII字符（包括中文）
                length += 2
            else:
                length += 1
        return length
    
    def _analyze_paragraph(self, para) -> Optional[ContentBlock]:
        """分析單個段落"""
        text = para.text.strip()
        if not text:
            return None
            
        formatting = self._extract_formatting(para)
        level, content_type = self._classify_content(text, formatting)
        
        return ContentBlock(
            text=text,
            level=level,
            content_type=content_type,
            formatting=formatting
        )

    def _classify_header(self, para) -> Optional[ContentBlock]:
        text = para.text.strip()
        if not text:
            return None

        formatting = self._extract_formatting(para)

        return ContentBlock(
            text=text,
            level=0,
            content_type='header',
            formatting=formatting
        )

    def _extract_formatting(self, para) -> Dict:
        """提取段落格式"""
        formatting = {
            'bold': False,
            'italic': False,
            'font_size': 12,
            'alignment': 'left'
        }
        
        if para.runs:
            run = para.runs[0]
            if run.bold:
                formatting['bold'] = True
            if run.italic:
                formatting['italic'] = True
            if run.font.size:
                formatting['font_size'] = run.font.size.pt
                
        return formatting
    
    def _classify_content(self, text: str, formatting: Dict) -> Tuple[int, str]:
        """分類內容類型和層級"""
        for pattern in self.chapter_patterns:
            if re.match(pattern, text):
                return (0, 'chapter')
        
        for pattern in self.subtitle_patterns:
            if re.match(pattern, text):
                return (1, 'subtitle')
            
        return (2, 'content')
    
    def _optimize_structure(self, blocks: List[ContentBlock]) -> List[ContentBlock]:
        """最佳化文件結構"""
        if not blocks:
            return blocks
            
        has_chapter = any(block.content_type == 'chapter' for block in blocks)
        if not has_chapter and blocks:
            for block in blocks:
                if block.level <= 1 or block.formatting.get('bold', False):
                    block.level = 0
                    block.content_type = 'chapter'
                    break
        
        return blocks

class ContentToSlideMapper:
    """內容到投影片的智慧映射器"""
    
    def __init__(self, template_path: str = None, template_content: bytes = None):
        self.template_path = template_path
        self.template_content = template_content
        self.prs = None
        self.max_content_length = 220  # 單張投影片最大內容長度
        self.max_content_items = 4     # 單張投影片最大內容項目數
        
    def create_slides(self, blocks: List[ContentBlock]) -> Presentation:
        """建立投影片（支援智慧分頁）- 清除範本原始內容"""
        if self.template_content:
            # 從bytes內容載入範本（用於Streamlit上傳的檔案）
            self.prs = Presentation(BytesIO(self.template_content))
        elif self.template_path:
            # 從檔案路徑載入範本
            self.prs = Presentation(self.template_path)
        else:
            raise ValueError("必須提供 template_path 或 template_content")
        
        # 清空現有投影片
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]
        
        current_slide = None
        current_content = []
        current_chapter_title = ""
        current_content_length = 0
        
        for block in blocks:
            print(f'Block level:{block.level}, type:{block.content_type}, text:{block.text[:50]}...')
            
            if block.content_type == 'header':
                # 完成當前投影片
                if current_slide is not None:
                    self._finalize_slide(current_slide, current_content)
                
                # 建立新的章節投影片
                current_slide = self._create_title_slide(block)
                current_content = []
                current_chapter_title = self._clean_chapter_text(block.text)
                current_content_length = 0
            
            elif block.content_type == 'chapter':
                # 完成當前投影片
                if current_slide is not None:
                    self._finalize_slide(current_slide, current_content)
                
                # 建立新的章節投影片
                current_slide = self._create_content_slide(block)
                current_content = []
                current_chapter_title = block.text
                current_content_length = 0

            elif block.content_type == 'subtitle':
                # 檢查是否需要分頁
                if True:  # 簡化邏輯，遇到subtitle始終分頁
                    # 完成當前投影片
                    if current_slide is not None:
                        self._finalize_slide(current_slide, current_content)
                    
                    # 建立新的內容投影片（使用相同章節標題）
                    chapter_block = ContentBlock(
                        text=current_chapter_title + " (續)",
                        level=0,
                        content_type='chapter',
                        formatting={}
                    )
                    current_slide = self._create_content_slide(chapter_block)
                    current_content = []
                    current_content_length = 0

                    current_content.append(block)
                    current_content_length += block.estimated_length
            else:
                # 檢查是否需要分頁
                if self._should_create_new_slide(current_content, current_content_length, block):
                    # 完成當前投影片
                    if current_slide is not None:
                        self._finalize_slide(current_slide, current_content)
                    
                    # 建立新的內容投影片（使用相同章節標題）
                    chapter_block = ContentBlock(
                        text=current_chapter_title + " (續)",
                        level=0,
                        content_type='chapter',
                        formatting={}
                    )
                    current_slide = self._create_content_slide(chapter_block)
                    current_content = []
                    current_content_length = 0
                
                # 添加內容到當前投影片
                if current_slide is None:
                    # 如果沒有當前投影片，建立一個
                    current_slide = self._create_content_slide(block)
                
                current_content.append(block)
                current_content_length += block.estimated_length
        
        # 處理最後一張投影片
        if current_slide is not None:
            self._finalize_slide(current_slide, current_content)
        
        print(f"已建立 {len(self.prs.slides)} 張新投影片")
        return self.prs
    
    def _should_create_new_slide(self, current_content: List[ContentBlock], 
                                current_length: int, new_block: ContentBlock) -> bool:
        """判斷是否需要建立新投影片"""
        # 檢查內容項目數量
        if len(current_content) >= self.max_content_items:
            return True
        
        # 檢查內容總長度
        if current_length + new_block.estimated_length > self.max_content_length:
            return True
        
        # 檢查是否為次標題（可能需要新投影片）
        if new_block.content_type == 'subtitle' and len(current_content) > 0:
            return True
        
        return False
    
    def _create_title_slide(self, block: ContentBlock):
        """建立標題投影片"""
        layout = self._get_best_layout('title')
        slide = self.prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            title_text = self._clean_chapter_text(block.text)
            slide.shapes.title.text = title_text
            
        return slide
    
    def _create_content_slide(self, block: ContentBlock):
        """建立內容投影片"""
        layout = self._get_best_layout('content')
        slide = self.prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            title_text = self._clean_subtitle_text(block.text)
            slide.shapes.title.text = title_text
            
        return slide
    
    def _clean_chapter_text(self, text: str) -> str:
        """清理章節文字，移除編號標記"""
        patterns = [
            r'^[一二三四五六七八九十]+[、．.]\s*',
            r'^第[一二三四五六七八九十壹貳參肆伍陸柒捌玖拾]+[章節部分]\s*',
            r'^第[一二三四五六七八九十]+[、．.]\s*',
            r'^[1-9]\d*[、．.]\s*',
            r'^第[1-9]\d*[章節部分]\s*',
            r'^第[1-9]\d*[、．.]\s*',
            r'^[A-Z][、．.]\s*',
            r'^第[A-Z][章節部分]\s*',
            r'^[●◆■▲]\s*',
        ]
        
        cleaned_text = text
        for pattern in patterns:
            cleaned_text = re.sub(pattern, '', cleaned_text)
        
        return cleaned_text.strip()
    
    def _clean_subtitle_text(self, text: str) -> str:
        """清理次標題文字"""
        patterns = [
            r'^[一二三四五六七八九十]+[）)]\s*',
            r'^\([一二三四五六七八九十]+\)\s*',
            r'^[1-9]\d*[）)]\s*',
            r'^\([1-9]\d*\)\s*',
            r'^[a-z][）)]\s*',
            r'^\([a-z]\)\s*',
            r'^[•·○]\s*',
        ]
        
        cleaned_text = text
        for pattern in patterns:
            cleaned_text = re.sub(pattern, '', cleaned_text)
        
        return cleaned_text.strip()
    
    def _finalize_slide(self, slide, content_blocks: List[ContentBlock]):
        """完成投影片內容"""
        if not content_blocks:
            return
            
        content_placeholder = None
        for shape in slide.placeholders:
            if hasattr(shape, 'text_frame') and shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break
        
        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            for i, block in enumerate(content_blocks):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = block.text
    
    def _get_best_layout(self, layout_type: str):
        """取得最佳佈局"""
        if not self.prs:
            raise Exception("Presentation not initialized")
            
        if layout_type == 'title':
            for layout in self.prs.slide_layouts:
                if 'title' in layout.name.lower() or layout == self.prs.slide_layouts[0]:
                    return layout
        else:
            for layout in self.prs.slide_layouts:
                if 'content' in layout.name.lower():
                    return layout
            if len(self.prs.slide_layouts) > 1:
                return self.prs.slide_layouts[1]
        
        return self.prs.slide_layouts[0]

class PPTXImageExporter:
    """PPTX 投影片圖片匯出器"""
    
    def __init__(self, output_dir: str = None):
        self.font_manager = SystemFontManager()
        if output_dir:
            self.temp_dir = output_dir
            os.makedirs(self.temp_dir, exist_ok=True)
        else:
            self.temp_dir = tempfile.mkdtemp()
        print(f"建立預覽圖片目錄: {self.temp_dir}")
    
    def export_slides_to_images(self, presentation_path: str = None, 
                               presentation_content: bytes = None) -> List[str]:
        """將 PPTX 投影片匯出為圖片"""
        try:
            print(f"使用改進的 python-pptx + PIL 方法處理投影片...")
            return self._export_with_enhanced_python_pptx(presentation_path, presentation_content)
                
        except Exception as e:
            print(f"圖片匯出失敗: {e}")
            return []
    
    def _export_with_enhanced_python_pptx(self, presentation_path: str = None, 
                                         presentation_content: bytes = None) -> List[str]:
        """使用改進的 python-pptx + PIL 匯出方法"""
        try:
            if presentation_content:
                prs = Presentation(BytesIO(presentation_content))
            elif presentation_path:
                prs = Presentation(presentation_path)
            else:
                raise ValueError("必須提供 presentation_path 或 presentation_content")
            
            image_paths = []
            
            print(f"使用改進的 python-pptx 方法處理 {len(prs.slides)} 張投影片...")
            
            for i, slide in enumerate(prs.slides, 1):
                print(f"正在處理投影片 {i}...")
                image_path = self._render_slide_to_image_enhanced(slide, i)
                if image_path:
                    image_paths.append(image_path)
            
            print(f"成功產生 {len(image_paths)} 張投影片圖片")
            return image_paths
            
        except Exception as e:
            print(f"改進的 python-pptx 匯出失敗: {e}")
            raise e
    
    def _render_slide_to_image_enhanced(self, slide, slide_number: int) -> str:
        """將投影片渲染為圖片（改進版）"""
        try:
            # 設定更高的解析度和更好的品質
            width, height = 1920, 1080  # 1080p 解析度
            img = Image.new('RGB', (width, height), color='white')
            draw = ImageDraw.Draw(img)
            
            # 使用跨平台字體管理器，更大的字體
            title_font = self.font_manager.get_best_font(64)
            content_font = self.font_manager.get_best_font(40)
            small_font = self.font_manager.get_best_font(28)
            
            # 繪製漸層背景
            self._draw_gradient_background(img, draw, width, height)
            
            # 提取和繪製內容
            y_pos = 80
            title_drawn = False
            content_items = []
            
            # 收集所有文字內容
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text and shape.text.strip():
                    raw_text = shape.text.strip()
                    text = self._normalize_text_cross_platform(raw_text)
                    
                    if not text:
                        continue
                    
                    is_title = (hasattr(shape, 'placeholder_format') and 
                              shape.placeholder_format.idx == 0)
                    
                    if is_title and not title_drawn:
                        # 繪製標題
                        y_pos = self._draw_title_enhanced_v2(draw, text, width, y_pos, title_font)
                        y_pos += 100
                        title_drawn = True
                    else:
                        content_items.append(text)
            
            # 繪製內容項目
            if content_items:
                y_pos = self._draw_content_enhanced_v2(draw, content_items, width, y_pos, content_font)
            
            # 添加投影片編號和裝飾
            self._add_slide_decorations_enhanced(draw, slide_number, width, height, small_font)
            
            # 應用後處理效果
            img = self._apply_post_processing(img)
            
            # 儲存高品質圖片
            image_path = os.path.join(self.temp_dir, f"slide_{slide_number}.jpg")
            img.save(image_path, 'JPEG', quality=95, optimize=True)
            
            print(f"成功建立改進版投影片圖片: {image_path}")
            return image_path
            
        except Exception as e:
            print(f"渲染投影片圖片時發生錯誤: {e}")
            return None
    
    def _draw_gradient_background(self, img: Image.Image, draw: ImageDraw.Draw, width: int, height: int):
        """繪製漸層背景"""
        try:
            # 創建垂直漸層
            for y in range(height):
                progress = y / height
                # 從淺藍到白色的漸層
                r = int(245 + (255 - 245) * progress)
                g = int(250 + (255 - 250) * progress)
                b = int(255)
                draw.line([(0, y), (width, y)], fill=(r, g, b))
        except Exception as e:
            print(f"繪製漸層背景時發生錯誤: {e}")
    
    def _draw_title_enhanced_v2(self, draw: ImageDraw.Draw, title: str, width: int, y_pos: int, font) -> int:
        """繪製增強版標題 v2"""
        try:
            safe_title = self._normalize_text_cross_platform(title)
            
            # 計算標題尺寸
            bbox = draw.textbbox((0, 0), safe_title, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
            
            # 置中位置
            center_x = (width - text_width) // 2
            
            # 繪製標題陰影（多層）
            for offset in range(1, 4):
                draw.text((center_x + offset, y_pos + offset), 
                         safe_title, fill=(0, 0, 0, 40), font=font)
            
            # 繪製主標題（使用漸層效果的顏色）
            draw.text((center_x, y_pos), safe_title, fill='#1A237E', font=font)
            
            # 繪製裝飾線條
            line_y = y_pos + text_height + 25
            line_thickness = 6
            
            # 主線條
            draw.rectangle([center_x, line_y, center_x + text_width, line_y + line_thickness], 
                          fill='#3498DB')
            
            # 裝飾點
            dot_size = 12
            for i in range(3):
                dot_x = center_x + (text_width // 4) * (i + 1)
                draw.ellipse([dot_x - dot_size//2, line_y + line_thickness + 10,
                             dot_x + dot_size//2, line_y + line_thickness + 10 + dot_size],
                            fill='#E74C3C')
            
            return line_y + line_thickness + 30
            
        except Exception as e:
            print(f"繪製標題時發生錯誤: {e}")
            return y_pos + 120
    
    def _draw_content_enhanced_v2(self, draw: ImageDraw.Draw, content_items: List[str], 
                                 width: int, y_pos: int, font) -> int:
        """繪製增強版內容 v2"""
        try:
            max_items = 8  # 增加最多顯示項目
            item_height = 80
            left_margin = 120
            bullet_size = 12
            
            for i, item in enumerate(content_items[:max_items]):
                if y_pos > 880:  # 1080p 高度限制，避免超出邊界
                    break
                
                safe_item = self._normalize_text_cross_platform(item)
                
                # 繪製項目符號（更精美的設計）
                bullet_x = left_margin - 40
                bullet_y = y_pos + 25
                
                # 外圈
                draw.ellipse([bullet_x - bullet_size, bullet_y - bullet_size,
                             bullet_x + bullet_size, bullet_y + bullet_size], 
                            fill='#3498DB', outline='#2980B9', width=2)
                
                # 內圈
                inner_size = bullet_size - 4
                draw.ellipse([bullet_x - inner_size, bullet_y - inner_size,
                             bullet_x + inner_size, bullet_y + inner_size], 
                            fill='white')
                
                # 項目背景（交替顏色）
                bg_color = '#F8F9FA' if i % 2 == 0 else '#FFFFFF'
                draw.rectangle([left_margin - 20, y_pos - 10,
                               width - 120, y_pos + item_height - 20],
                              fill=bg_color, outline='#E9ECEF', width=1)
                
                # 處理長文字換行（改進版）
                wrapped_lines = self._wrap_text_smart_v2(safe_item, width - left_margin - 140, font, draw)
                
                line_y = y_pos
                for j, line in enumerate(wrapped_lines[:3]):  # 最多3行
                    if line.strip():
                        # 第一行使用較深的顏色
                        text_color = '#2C3E50' if j == 0 else '#5D6D7E'
                        draw.text((left_margin, line_y), line, fill=text_color, font=font)
                        line_y += 42
                
                y_pos += item_height
            
            return y_pos
            
        except Exception as e:
            print(f"繪製內容時發生錯誤: {e}")
            return y_pos + 300
    
    def _add_slide_decorations_enhanced(self, draw: ImageDraw.Draw, slide_number: int, 
                                       width: int, height: int, font):
        """添加增強版投影片裝飾元素"""
        try:
            # 頁腳背景
            footer_height = 60
            draw.rectangle([0, height - footer_height, width, height], 
                          fill='#34495E', outline='#2C3E50', width=2)
            
            # 投影片編號（左下角）
            number_text = f"第 {slide_number} 頁"
            draw.text((30, height - 45), number_text, fill='white', font=font)
            
            # 系統資訊（右下角）
            system_text = f"Python-PPTX Enhanced Renderer | {platform.system()}"
            bbox = draw.textbbox((0, 0), system_text, font=font)
            text_width = bbox[2] - bbox[0]
            draw.text((width - text_width - 30, height - 45), system_text, 
                     fill='#BDC3C7', font=font)
            
            # 頂部裝飾帶
            draw.rectangle([0, 0, width, 8], fill='#3498DB')
            draw.rectangle([0, 8, width, 16], fill='#E74C3C')
            
            # 右側裝飾元素
            for i in range(5):
                y = 200 + i * 150
                draw.rectangle([width - 20, y, width - 10, y + 80], 
                              fill='#95A5A6', outline='#7F8C8D', width=1)
            
        except Exception as e:
            print(f"添加裝飾時發生錯誤: {e}")
    
    def _wrap_text_smart_v2(self, text: str, max_width: int, font, draw) -> List[str]:
        """智慧文字換行 v2"""
        try:
            lines = []
            words = text.split()
            current_line = ""
            
            for word in words:
                test_line = current_line + (" " if current_line else "") + word
                
                try:
                    bbox = draw.textbbox((0, 0), test_line, font=font)
                    text_width = bbox[2] - bbox[0]
                    
                    if text_width <= max_width:
                        current_line = test_line
                    else:
                        if current_line:
                            lines.append(current_line)
                        # 如果單個詞太長，截斷它
                        if len(word) > max_width // 30:
                            word = word[:max_width//30] + "..."
                        current_line = word
                except:
                    # 備用方法：按字符數估算
                    if len(test_line) < max_width // 25:
                        current_line = test_line
                    else:
                        if current_line:
                            lines.append(current_line)
                        current_line = word
            
            if current_line:
                lines.append(current_line)
            
            return lines
            
        except Exception as e:
            print(f"智慧文字換行處理錯誤: {e}")
            return [text[:60] + "..." if len(text) > 60 else text]
    
    def _apply_post_processing(self, img: Image.Image) -> Image.Image:
        """應用後處理效果"""
        try:
            # 輕微銳化
            img = img.filter(ImageFilter.UnsharpMask(radius=1, percent=120, threshold=2))
            
            # 增強對比度
            enhancer = ImageEnhance.Contrast(img)
            img = enhancer.enhance(1.1)
            
            # 增強顏色飽和度
            enhancer = ImageEnhance.Color(img)
            img = enhancer.enhance(1.05)
            
            return img
            
        except Exception as e:
            print(f"後處理時發生錯誤: {e}")
            return img
    
    def _normalize_text_cross_platform(self, text: str) -> str:
        """跨平台標準化文字編碼"""
        try:
            if isinstance(text, bytes):
                # 依平台嘗試不同編碼
                if platform.system() == "Darwin":  # macOS
                    encodings = ['utf-8', 'utf-8-sig', 'macroman', 'big5']
                elif platform.system() == "Windows":
                    encodings = ['utf-8', 'utf-8-sig', 'cp950', 'big5', 'gb2312']
                else:
                    encodings = ['utf-8', 'utf-8-sig', 'gb2312', 'big5']
                
                for encoding in encodings:
                    try:
                        text = text.decode(encoding)
                        break
                    except:
                        continue
                else:
                    text = text.decode('utf-8', errors='ignore')
            
            text = str(text)
            text = text.encode('utf-8', errors='ignore').decode('utf-8')
            
            # 移除控制字元但保留換行符
            import unicodedata
            text = ''.join(char for char in text 
                          if unicodedata.category(char)[0] != 'C' or char in '\n\r\t')
            
            return text.strip()
            
        except Exception as e:
            print(f"跨平台文字編碼標準化失敗: {e}")
            return str(text)[:100] if text else "文字顯示錯誤"
    
    def cleanup(self, force_delete: bool = False):
        """清理暫存檔案"""
        try:
            import shutil
            if os.path.exists(self.temp_dir) and force_delete:
                shutil.rmtree(self.temp_dir)
                print(f"清理目錄: {self.temp_dir}")
        except Exception as e:
            print(f"清理檔案時發生錯誤: {e}")

class WordToPPTXConverter:
    """Word轉PowerPoint轉換器主類別"""
    
    def __init__(self):
        self.analyzer = WordDocumentAnalyzer()
        self.font_manager = SystemFontManager()
    
    def convert(self, word_file_path: str = None, word_file_content: bytes = None,
                template_file_path: str = None, template_file_content: bytes = None,
                output_path: str = None, 
                generate_preview: bool = True,
                save_preview_to_disk: bool = True) -> ConversionResult:
        """
        轉換Word文件為PowerPoint
        
        Args:
            word_file_path: Word檔案路徑（獨立應用程式使用）
            word_file_content: Word檔案內容（Streamlit使用）
            template_file_path: 範本檔案路徑（獨立應用程式使用）
            template_file_content: 範本檔案內容（Streamlit使用）
            output_path: 輸出檔案路徑
            generate_preview: 是否生成預覽圖片
            save_preview_to_disk: 是否保存預覽圖片到磁碟（預設True）
            
        Returns:
            ConversionResult: 轉換結果
        """
        try:
            # 步驟1：分析Word文件
            print("正在分析Word文件章節結構...")
            blocks = self.analyzer.analyze_document(
                file_path=word_file_path, 
                file_content=word_file_content
            )
            
            if not blocks:
                return ConversionResult(
                    success=False,
                    error_message="Word文件中沒有找到可轉換的內容"
                )
            
            # 步驟2：建立投影片
            print("正在建立PowerPoint投影片...")
            mapper = ContentToSlideMapper(
                template_path=template_file_path,
                template_content=template_file_content
            )
            presentation = mapper.create_slides(blocks)
            
            # 步驟3：儲存檔案
            print("正在儲存PowerPoint檔案...")
            if output_path:
                output_dir = os.path.dirname(output_path)
                if output_dir and not os.path.exists(output_dir):
                    os.makedirs(output_dir, exist_ok=True)
                presentation.save(output_path)
            
            # 步驟4：生成預覽圖片（可選）
            preview_images = []
            if generate_preview:
                print("正在生成預覽圖片...")
                
                if save_preview_to_disk and output_path:
                    # 保存到磁碟模式：在輸出檔案同目錄創建預覽資料夾
                    output_dir = os.path.dirname(output_path)
                    output_name = os.path.splitext(os.path.basename(output_path))[0]
                    preview_dir = os.path.join(output_dir, f"{output_name}_預覽圖片")
                    exporter = PPTXImageExporter(preview_dir)
                    
                    # 從磁碟檔案生成預覽
                    preview_images = exporter.export_slides_to_images(output_path)
                    
                elif save_preview_to_disk:
                    # 保存到磁碟模式但沒有指定輸出路徑：使用暫存目錄
                    preview_dir = tempfile.mkdtemp()
                    exporter = PPTXImageExporter(preview_dir)
                    
                    # 從記憶體中的presentation生成預覽
                    temp_pptx_path = os.path.join(tempfile.mkdtemp(), "temp_presentation.pptx")
                    presentation.save(temp_pptx_path)
                    preview_images = exporter.export_slides_to_images(temp_pptx_path)
                    
                    # 清理暫存的pptx檔案
                    try:
                        os.remove(temp_pptx_path)
                        os.rmdir(os.path.dirname(temp_pptx_path))
                    except:
                        pass
                        
                else:
                    # 僅記憶體模式：生成預覽但不保存到磁碟
                    print("生成記憶體預覽模式...")
                    preview_dir = tempfile.mkdtemp()
                    exporter = PPTXImageExporter(preview_dir)
                    
                    # 從記憶體中的presentation生成預覽
                    temp_pptx_path = os.path.join(tempfile.mkdtemp(), "temp_presentation.pptx")
                    presentation.save(temp_pptx_path)
                    preview_images = exporter.export_slides_to_images(temp_pptx_path)
                    
                    # 清理暫存檔案（因為不保存到磁碟）
                    try:
                        os.remove(temp_pptx_path)
                        os.rmdir(os.path.dirname(temp_pptx_path))
                        # 清理預覽圖片目錄（僅記憶體模式）
                        exporter.cleanup(force_delete=True)
                    except:
                        pass
            
            return ConversionResult(
                success=True,
                output_path=output_path or "",
                preview_images=preview_images,
                slides_count=len(presentation.slides)
            )
            
        except Exception as e:
            import traceback
            error_details = f"{str(e)}\n\n詳細錯誤資訊:\n{traceback.format_exc()}"
            return ConversionResult(
                success=False,
                error_message=error_details
            )

# 驗證相依性
def check_dependencies():
    """檢查相依性"""
    return DEPENDENCIES_OK

def get_dependency_status():
    """取得相依性狀態"""
    status = {}
    try:
        import docx
        status['python-docx'] = True
    except ImportError:
        status['python-docx'] = False
    
    try:
        import pptx
        status['python-pptx'] = True
    except ImportError:
        status['python-pptx'] = False
    
    try:
        import PIL
        status['pillow'] = True
    except ImportError:
        status['pillow'] = False
    
    return status